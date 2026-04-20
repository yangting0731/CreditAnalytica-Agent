"""PPT 导出模块 — 将分析报告导出为 PowerPoint"""
import io
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import plotly.graph_objects as go


def _add_title_slide(prs: Presentation, title: str, subtitle: str):
    """添加封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)


def _add_section_slide(prs: Presentation, section_title: str, fig: go.Figure,
                        insight_text: str):
    """添加分析维度页 — 图表 + 结论"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)
    # Insight text
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = insight_text
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
    # Chart image
    if fig is not None:
        try:
            img_bytes = fig.to_image(format="png", width=800, height=400, scale=2)
            img_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(img_stream, Inches(0.5), Inches(2.2), Inches(9), Inches(4.5))
        except Exception:
            # Fallback if kaleido not available
            txBox3 = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(1))
            tf3 = txBox3.text_frame
            p3 = tf3.paragraphs[0]
            p3.text = "[图表导出需要安装 kaleido 库]"
            p3.font.size = Pt(14)
            p3.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
            p3.alignment = PP_ALIGN.CENTER


def _add_strategy_slide(prs: Presentation, strategy_text: str):
    """添加策略建议页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "策略建议"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)
    # Strategy content
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    # Split by sections
    for line in strategy_text.split("\n"):
        line = line.strip()
        if not line:
            continue
        p2 = tf2.add_paragraph()
        if line.startswith("##"):
            p2.text = line.replace("##", "").strip()
            p2.font.size = Pt(14)
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(0x25, 0x63, 0xEB)
            p2.space_before = Pt(12)
        else:
            p2.text = line
            p2.font.size = Pt(10)
            p2.font.color.rgb = RGBColor(0x37, 0x41, 0x51)


def generate_pptx(sections: list, strategy_text: str, title: str = "客群分析报告") -> bytes:
    """
    生成 PPT 文件

    Args:
        sections: list of dict, each with keys:
            - title: str (section title)
            - fig: go.Figure or None
            - insight: str (analysis conclusion)
        strategy_text: str (strategy report text)
        title: str (report title)

    Returns:
        bytes: PPT file content
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Cover slide
    _add_title_slide(prs, title, "百融智能分析 Agent 生成")

    # Section slides
    for section in sections:
        _add_section_slide(prs, section["title"], section.get("fig"), section.get("insight", ""))

    # Strategy slide
    if strategy_text:
        _add_strategy_slide(prs, strategy_text)

    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
